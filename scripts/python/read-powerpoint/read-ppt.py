#!/usr/local/Caskroom/miniconda/base/bin/python

import os
import pptx
import configparser
import requests


def process_file(folder, api_key):
    print('traversing {}'.format(folder))
    if 'clean' in folder:
        return
    for filename in os.listdir(folder):
        file_path = os.path.join(folder, filename)
        if filename.endswith(".pptx"):
            text = get_text_from_ppt(os.path.join(file_path))
            text = sanitize_string(text)
            # print("sanitized string: {}".format(text))
            ai_response = send_to_open_ai(text, api_key)
            print(ai_response)
            print("======== done ===========")
            continue;
        elif os.path.isfile(file_path):
            continue
        else:
            process_file(file_path, api_key)


def send_to_open_ai(text, api_key):
    prompt = 'Give me a synthesize paragraph from the following text: {}'.format(text)
    data = {'prompt': prompt, 'max_tokens': 500}
    headers = {'Authorization': 'Bearer {}'.format(api_key)}
    response = requests.post('https://api.openai.com/v1/engines/text-davinci-003/completions', json=data,
                             headers=headers)
    if (response.status_code == 400):
        return response.text

    raw_response = response.json()
    ai_response = raw_response["choices"][0]["text"]
    return ai_response


def get_text_from_ppt(ppt_path):
    prs = pptx.Presentation(ppt_path)
    print("===== processing file: {} =======".format(ppt_path))
    text = []
    for slide in prs.slides:
        slide_id = prs.slides.index(slide)
        text.append("Title: {} slide: {}".format(ppt_path, slide_id))
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
    return '\n'.join(text)


def sanitize_string(s):
    lines = s.split('\n')
    new_lines = []
    for line in lines:
        if len(line.split(' ')) <= 3:
            continue
        if line in new_lines:
            continue
        if line == '':
            continue
        if line.startswith('Title:'):
            continue
        new_lines.append(line)
    return '\n'.join(new_lines)


def main():
    config = configparser.ConfigParser()
    config.read(os.path.join(os.path.dirname(__file__), 'properties.ini'))
    directory_ = config['DEFAULT']['Local_Directory']
    print("reading: {}".format(directory_))
    process_file(directory_, config['DEFAULT']['API_KEY'])


# if __name__openai_key = "YOUR_OPENAI_KEY_HERE"

if __name__ == "__main__":
    main()
